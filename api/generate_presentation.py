import pathlib
import textwrap
import google.generativeai as genai
from IPython.display import display
from IPython.display import Markdown
from pptx import Presentation
from pptx.util import Inches
import textwrap
import requests
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from pptx.enum.text import PP_ALIGN
from flask import Flask, render_template, request

PEXELS_API_KEY= 'nSV0ObrWctMxhpxeFI38kh8mAblRdnjha1cQ50OaMoI3nCbw1QPHZjP4'
GOOGLE_API_KEY='AIzaSyCFqqKkv-LEps4-ejViZWzqZEJ6dKsxooA'
genai.configure(api_key=GOOGLE_API_KEY)
# userdata.get('GOOGLE_API_KEY')
model = genai.GenerativeModel('gemini-1.5-flash')

app = Flask(__name__)

def to_markdown(text):
  text = text.replace('•', '  *')
  return Markdown(textwrap.indent(text, '> ', predicate=lambda _: True))

def fetch_image_url(query):
    url = 'https://api.pexels.com/v1/search'
    headers = {
        'Authorization': PEXELS_API_KEY
    }
    params = {
        'query': query,
        'per_page': 1
    }
    response = requests.get(url, headers=headers, params=params)
    if response.status_code == 200:
        data = response.json()
        if data['photos']:
            return data['photos'][0]['src']['large']
    return None

@app.route('/api/generate_presentation', methods=['GET'])
def generate_presentation():
    title=request.args.get("title")
    grade=request.args.get("grade")
    num_slides=request.args.get("num_slides")

    # if title is None or grade is None or num_slides is None:
    #     return render_template("index.html")

    response = model.generate_content(f"""
        Give content for {num_slides} slides on the topic '{title}' for {grade} Students
        in the following format such as dont repeat title i given to you just start writing content:
        'Short Title of Slide and dont write Slide No. and also dont use ##'
        • Content bullet points
        Note: Slides Format should be strictly like:
        **The Big Picture**

        * Quick Sort is a highly efficient sorting algorithm.
        * It works by dividing the list into sub-lists and recursively sorting them.
        * It is known for its average-case time complexity of O(n log n).
        """)
    to_markdown(response.text)
    content=str(response.text)
    print(content)

    # Split the content by new lines and filter out empty lines
    lines = [line for line in content.split('\n') if line.strip()]

    # Initialize lists to hold the slide titles and content
    slide_titles = []
    slide_contents = []
    current_content = []

    # Flag to start collecting content after the first ##
    collecting = False

    # Iterate over each line
    for line in lines:

        # Check if the line starts with '**' indicating a new slide title
        if line.startswith('## **'):
            # If we have collected content for a previous slide, add it to the list
            if current_content:
                slide_contents.append("\n".join(current_content))
                current_content = []
            # Extract the title and add it to the list of titles
            slide_titles.append(line.split("**")[1].strip())
        elif line.startswith('**'):
            # If we have collected content for a previous slide, add it to the list
            if current_content:
                slide_contents.append("\n".join(current_content))
                current_content = []
            # Extract the title and add it to the list of titles
            slide_titles.append(line.split("**")[1].strip())
        else:
            # Add the line to the current slide's content
            cleaned_line = line.lstrip('* ').strip()
            current_content.append(cleaned_line)

    # Don't forget to add the last slide's content
    if current_content:
        slide_contents.append("\n".join(current_content))

    # Create a presentation object
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Create slides with titles and contents
    for title, content in zip(slide_titles, slide_contents):
        slide_layout = prs.slide_layouts[1]  # Use the 'Title and Content' layout
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = title
        for paragraph in title_placeholder.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

        content_placeholder.text = content

        # Fetch the image URL
        image_url = fetch_image_url(title)
        if image_url:
            response = requests.get(image_url)
            if response.status_code == 200:
                image_stream = BytesIO(response.content)
                left = Inches(8.8)
                top = Inches(1.5)
                slide.shapes.add_picture(image_stream, left, top, width=Inches(4.5), height=Inches(6))

    # Save the presentation
    prs.save('presentation.pptx')

    print(len(slide_titles))
    print(len(slide_contents))

    print("PowerPoint presentation created successfully!")
    return render_template("index.html")

if __name__ == "__main__":
    app.run(host="192.168.56.1", port=9000, debug=True)